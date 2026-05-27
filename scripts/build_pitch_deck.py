from __future__ import annotations

import argparse
import re
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path
from tempfile import NamedTemporaryFile
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "docs" / "sati_pitch.pptx"

BG = "FAF8F3"
GREEN = "6FB866"
WARM = "DCA068"
TEXT = "34492F"
MUTED = "6B7566"
CREAM = "FFFDF8"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


@dataclass(frozen=True)
class SlideSpec:
    title: str
    body: list[str]
    note: str
    kind: str = "bullets"
    kicker: str = ""


SLIDES = [
    SlideSpec(
        title="Sati",
        kicker="Posture & Focus Coach",
        body=[
            "Hackathon Coding Thailand 2026 · Wellness Track",
            "A sensor-driven wellness companion that helps office workers notice posture",
            "and break patterns before they drift too far.",
        ],
        note=(
            "Sati — Posture & Focus Coach\n"
            "Hackathon Coding Thailand 2026 · Wellness Track\n\n"
            'Key line: "A sensor-driven wellness companion that helps office workers notice posture and break patterns before they drift too far."'
        ),
        kind="title",
    ),
    SlideSpec(
        title="Problem",
        body=[
            "Office workers sit in front of screens for long hours, and posture often changes without awareness.",
            "Existing habit apps rely on self-report, so users forget to log.",
            "Many wellness apps do not have hardware ground truth.",
        ],
        note=(
            "Slide 2: Problem (30s)\n\n"
            "- Office workers sit in front of screens for long hours, and posture often changes without awareness.\n"
            "- Existing habit apps rely on self-report, so users forget to log.\n"
            "- Many wellness apps do not have hardware ground truth.\n\n"
            'Speaker note:\n"The problem is not that people do not care. The problem is that the signal arrives too late, and most apps ask the user to do extra work."'
        ),
    ),
    SlideSpec(
        title="Insight",
        body=[
            "คนไม่กลับมาเปิดแอป ถ้ามันรู้สึกเหมือนการบ้าน",
            "Ambient feedback instead of interruption-first design",
            "Reward loop that grows from measured behavior",
            "Cute plant mechanic that makes progress visible",
        ],
        note=(
            "Slide 3: Insight (30s)\n\n"
            '"คนไม่กลับมาเปิดแอป ถ้ามันรู้สึกเหมือนการบ้าน"\n\n'
            "What Sati changes:\n"
            "- Ambient feedback instead of interruption-first design\n"
            "- Reward loop that grows from measured behavior\n"
            "- Cute plant mechanic that makes progress visible"
        ),
        kind="quote",
    ),
    SlideSpec(
        title="Solution",
        body=[
            "IMU on Nano 33 BLE Sense reads back angle.",
            "Modulino Distance reads screen distance.",
            "Movement cue confirms break behavior.",
            "Web app turns signals into a 3-state coach: NORMAL -> WARNING -> ACTION.",
            "Growth mechanic makes the plant grow from real sensor-confirmed behavior.",
            "Second-Brain view summarizes observed behavior patterns.",
            "Persona Avatar form creates an LLM-ready brief for avatar personalization.",
        ],
        note=(
            "Slide 4: Solution (45s)\n\n"
            "- IMU on Nano 33 BLE Sense reads back angle.\n"
            "- Modulino Distance reads screen distance.\n"
            "- Movement cue confirms break behavior.\n"
            "- Web app turns signals into a 3-state coach: NORMAL -> WARNING -> ACTION.\n"
            "- Growth mechanic makes the plant grow from real sensor-confirmed behavior.\n"
            "- Second-Brain view summarizes observed behavior patterns.\n"
            "- Persona Avatar form creates an LLM-ready brief for avatar personalization."
        ),
        kind="solution",
    ),
    SlideSpec(
        title="Architecture",
        body=[
            "[Nano 33 BLE Sense] --BLE--> [UNO Q Linux]",
            "[Modulino ToF]      --Serial-> [UNO Q MCU]",
            "[UNO Q Linux]       --WebSocket--> [Next.js Browser]",
            "[Browser]           --window.name--> demo memory",
        ],
        note=(
            "Slide 5: Architecture (30s)\n\n"
            "[Nano 33 BLE Sense] --BLE--> [Arduino UNO Q Linux] --WebSocket--> [Next.js Browser]\n"
            "[Modulino ToF]      --Serial-> [UNO Q MCU] ----------------------^\n\n"
            'Speaker note:\n"UNO Q is the edge hub. The MCU side handles realtime sensor bridge. The Linux side runs Python and serves the dashboard."'
        ),
        kind="architecture",
    ),
    SlideSpec(
        title="Demo (Live)",
        body=[
            "Normal posture",
            "Hunched cue",
            "Stretch completion",
            "Growth level-up",
            "Second-Brain insights",
            "Persona Avatar recommendation",
        ],
        note=(
            "Slide 6: Demo (Live)\n\n"
            "Follow docs/DEMO_SCRIPT.md:\n\n"
            "1. Normal posture\n"
            "2. Hunched cue\n"
            "3. Stretch completion\n"
            "4. Growth level-up\n"
            "5. Second-Brain insights\n"
            "6. Persona Avatar recommendation"
        ),
        kind="demo",
    ),
    SlideSpec(
        title="Market & Business",
        body=[
            "Beachhead: B2B office wellness pilot for tech companies and co-working spaces.",
            "Revenue: hardware kit + optional wellness insight reports.",
            "Expansion: decorative content, team challenges, privacy-preserving team-ready reports.",
            "Market sizing source to be added before final submission.",
        ],
        note=(
            "Slide 7: Market & Business (30s)\n\n"
            "- Beachhead: B2B office wellness pilot for tech companies and co-working spaces.\n"
            "- Revenue: hardware kit + optional wellness insight reports.\n"
            "- Expansion: decorative content, team challenges, privacy-preserving team-ready reports.\n"
            "- Market size: [ทีมต้องเติม source จากรายงาน corporate wellness / workplace wellness ที่น่าเชื่อถือ]\n\n"
            'Speaker note:\n"We start where the hardware value is obvious: offices that already invest in wellness but cannot see behavior patterns in realtime."'
        ),
    ),
    SlideSpec(
        title="Ask",
        body=[
            "Feedback from judges on hardware reliability and pilot design",
            "Mentorship on enterprise wellness sales",
            "Distribution partnership with Arduino, maker community, and workplace wellness partners",
            "Pilot company for a 2-week office trial",
        ],
        note=(
            "Slide 8: Ask\n\n"
            "- Feedback from judges on hardware reliability and pilot design\n"
            "- Mentorship on enterprise wellness sales\n"
            "- Distribution partnership with Arduino / maker community / corporate wellness partners\n"
            "- Pilot company for a 2-week office trial"
        ),
    ),
    SlideSpec(
        title="Closing",
        body=["sensor -> insight -> action -> reward -> daily awareness"],
        note=(
            "Closing Line\n\n"
            '"Sati is not just a reminder. It is a loop: sensor -> insight -> action -> reward -> daily awareness."'
        ),
        kind="closing",
    ),
]


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color(BG)
    bg.line.fill.background()


def set_text(run, size: int, value: str = TEXT, bold: bool = False, font: str = "Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(value)


def add_slide_number(slide, slide_no: int, total: int):
    box = slide.shapes.add_textbox(Inches(12.1), Inches(6.95), Inches(0.9), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{slide_no}/{total}"
    set_text(r, 11, MUTED)


def add_header(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.42), Inches(8.5), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_text(r, 32, TEXT, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.17), Inches(11.8), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color(GREEN)
    line.line.fill.background()


def add_bullets(slide, bullets: list[str], x=0.9, y=1.65, w=11.5, h=4.8, size=28, spacing=10):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = f"• {item}"
        set_text(r, size, TEXT)


def add_footer_mark(slide):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.78), Inches(2.1), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color("E9F4E3")
    pill.line.color.rgb = color("C9DFBF")
    box = slide.shapes.add_textbox(Inches(0.96), Inches(6.86), Inches(1.7), Inches(0.18))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Sati demo deck"
    set_text(r, 10, TEXT, bold=True)


def render_title(slide, spec: SlideSpec):
    add_background(slide)
    leaf = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(0.75), Inches(0.7), Inches(0.7))
    leaf.fill.solid()
    leaf.fill.fore_color.rgb = color(GREEN)
    leaf.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.05), Inches(1.8), Inches(11.2), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = spec.title
    set_text(r, 68, TEXT, bold=True)

    kicker = slide.shapes.add_textbox(Inches(1.05), Inches(3.0), Inches(11.2), Inches(0.55))
    p = kicker.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = spec.kicker
    set_text(r, 30, GREEN, bold=True)

    for i, line in enumerate(spec.body):
        box = slide.shapes.add_textbox(Inches(1.65), Inches(4.05 + i * 0.68), Inches(10.0), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        set_text(r, 20 if i == 0 else 18, TEXT if i == 0 else MUTED)


def render_quote(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)

    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.65), Inches(11.55), Inches(1.65))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = color(CREAM)
    quote_box.line.color.rgb = color("D8E7D0")

    text = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.95), Inches(0.9))
    p = text.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = spec.body[0]
    set_text(r, 28, TEXT, bold=True)

    add_bullets(slide, spec.body[1:], y=3.75, size=29, spacing=12)


def render_architecture(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.65), Inches(11.65), Inches(3.8))
    frame.fill.solid()
    frame.fill.fore_color.rgb = color(CREAM)
    frame.line.color.rgb = color("C9DFBF")

    diag = "\n".join(spec.body)
    box = slide.shapes.add_textbox(Inches(1.15), Inches(2.35), Inches(11.0), Inches(2.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = diag
    set_text(r, 22, TEXT, bold=True, font="Consolas")

    note = slide.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.6))
    p = note.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "UNO Q = edge hub: MCU for realtime bridge, Linux for Python WebSocket + static dashboard"
    set_text(r, 18, MUTED)


def render_solution(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)
    add_bullets(slide, spec.body, y=1.5, size=25, spacing=8)

    arch = slide.shapes.add_textbox(Inches(5.25), Inches(6.05), Inches(6.8), Inches(0.48))
    p = arch.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "Nano IMU + Modulino ToF + UNO Q bridge -> Browser feedback"
    set_text(r, 14, MUTED)


def render_demo(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)
    add_bullets(slide, spec.body, x=0.95, y=1.55, w=7.2, h=4.8, size=26, spacing=7)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(2.4), Inches(3.15), Inches(1.25))
    box.fill.solid()
    box.fill.fore_color.rgb = color("F7E1C5")
    box.line.color.rgb = color(WARM)
    text = slide.shapes.add_textbox(Inches(8.86), Inches(2.8), Inches(2.75), Inches(0.48))
    p = text.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Live Demo"
    set_text(r, 30, TEXT, bold=True)


def render_closing(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)
    box = slide.shapes.add_textbox(Inches(1.15), Inches(2.55), Inches(11.0), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = spec.body[0]
    set_text(r, 34, TEXT, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(4.45), Inches(7.3), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color(WARM)
    accent.line.fill.background()


def render_default(slide, spec: SlideSpec):
    add_background(slide)
    add_header(slide, spec.title)
    add_bullets(slide, spec.body)


def build_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if spec.kind == "title":
            render_title(slide, spec)
        elif spec.kind == "quote":
            render_quote(slide, spec)
        elif spec.kind == "architecture":
            render_architecture(slide, spec)
        elif spec.kind == "solution":
            render_solution(slide, spec)
        elif spec.kind == "demo":
            render_demo(slide, spec)
        elif spec.kind == "closing":
            render_closing(slide, spec)
        else:
            render_default(slide, spec)
        add_footer_mark(slide)
        add_slide_number(slide, idx, len(SLIDES))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    add_speaker_notes(output_path, [slide.note for slide in SLIDES])


def next_rid(xml: str) -> str:
    ids = [int(match.group(1)) for match in re.finditer(r'Id="rId(\d+)"', xml)]
    return f"rId{max(ids, default=0) + 1}"


def add_relationship(xml: str, rel_id: str, rel_type: str, target: str) -> str:
    if f'Target="{target}"' in xml:
        return xml
    rel = f'<Relationship Id="{rel_id}" Type="{rel_type}" Target="{target}"/>'
    return xml.replace("</Relationships>", f"{rel}</Relationships>", 1)


def add_content_type(xml: str, part_name: str, content_type: str) -> str:
    if f'PartName="{part_name}"' in xml:
        return xml
    override = f'<Override PartName="{part_name}" ContentType="{content_type}"/>'
    return xml.replace("</Types>", f"{override}</Types>", 1)


def add_notes_master_id(xml: str, rel_id: str) -> str:
    if "<p:notesMasterIdLst>" in xml:
        return xml
    fragment = f'<p:notesMasterIdLst><p:notesMasterId r:id="{rel_id}"/></p:notesMasterIdLst>'
    if "</p:sldMasterIdLst>" in xml:
        return xml.replace("</p:sldMasterIdLst>", f"</p:sldMasterIdLst>{fragment}", 1)
    return xml.replace("<p:sldIdLst>", f"{fragment}<p:sldIdLst>", 1)


def notes_master_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="685800"/><a:ext cx="5486400" cy="3086100"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="2"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="4114800"/><a:ext cx="7772400" cy="2926080"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" sz="1200"/></a:p></p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>
"""


def notes_paragraphs_xml(note: str) -> str:
    paragraphs = []
    for line in note.splitlines():
        if line.strip():
            paragraphs.append(
                '<a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>'
                f"{escape(line)}"
                "</a:t></a:r><a:endParaRPr lang=\"en-US\" sz=\"1200\"/></a:p>"
            )
        else:
            paragraphs.append('<a:p><a:endParaRPr lang="en-US" sz="1200"/></a:p>')
    return "".join(paragraphs)


def notes_slide_xml(note: str) -> str:
    body = notes_paragraphs_xml(note)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="685800"/><a:ext cx="5486400" cy="3086100"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="2"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="4114800"/><a:ext cx="7772400" cy="2926080"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/>{body}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>
"""


def notes_slide_rels(slide_number: int) -> str:
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_number}.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
</Relationships>
"""


def notes_master_rels() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>
"""


def add_speaker_notes(pptx_path: Path, notes: list[str]) -> None:
    with zipfile.ZipFile(pptx_path, "r") as src:
        presentation_rels = src.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        notes_master_rid = next_rid(presentation_rels)

        with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp_path = Path(tmp.name)

        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                data = src.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    xml = data.decode("utf-8")
                    xml = add_content_type(
                        xml,
                        "/ppt/notesMasters/notesMaster1.xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml",
                    )
                    for idx in range(1, len(notes) + 1):
                        xml = add_content_type(
                            xml,
                            f"/ppt/notesSlides/notesSlide{idx}.xml",
                            "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
                        )
                    data = xml.encode("utf-8")
                elif item.filename == "ppt/_rels/presentation.xml.rels":
                    xml = data.decode("utf-8")
                    xml = add_relationship(
                        xml,
                        notes_master_rid,
                        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
                        "notesMasters/notesMaster1.xml",
                    )
                    data = xml.encode("utf-8")
                elif item.filename == "ppt/presentation.xml":
                    xml = data.decode("utf-8")
                    xml = add_notes_master_id(xml, notes_master_rid)
                    data = xml.encode("utf-8")
                else:
                    match = re.fullmatch(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", item.filename)
                    if match:
                        slide_number = int(match.group(1))
                        if slide_number <= len(notes):
                            xml = data.decode("utf-8")
                            xml = add_relationship(
                                xml,
                                next_rid(xml),
                                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                                f"../notesSlides/notesSlide{slide_number}.xml",
                            )
                            data = xml.encode("utf-8")
                dst.writestr(item, data)

            dst.writestr("ppt/notesMasters/notesMaster1.xml", notes_master_xml())
            dst.writestr("ppt/notesMasters/_rels/notesMaster1.xml.rels", notes_master_rels())
            for idx, note in enumerate(notes, start=1):
                dst.writestr(f"ppt/notesSlides/notesSlide{idx}.xml", notes_slide_xml(note))
                dst.writestr(f"ppt/notesSlides/_rels/notesSlide{idx}.xml.rels", notes_slide_rels(idx))

    shutil.move(tmp_path, pptx_path)


def render_preview(preview_dir: Path, slide_numbers: list[int]) -> None:
    from PIL import Image, ImageDraw, ImageFont

    preview_dir.mkdir(parents=True, exist_ok=True)
    width, height = 1600, 900
    try:
        font_big = ImageFont.truetype("calibri.ttf", 86)
        font_title = ImageFont.truetype("calibri.ttf", 54)
        font_body = ImageFont.truetype("calibri.ttf", 36)
        font_mono = ImageFont.truetype("consola.ttf", 34)
    except OSError:
        font_big = ImageFont.load_default()
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()
        font_mono = ImageFont.load_default()

    for slide_number in slide_numbers:
        spec = SLIDES[slide_number - 1]
        img = Image.new("RGB", (width, height), f"#{BG}")
        draw = ImageDraw.Draw(img)
        if spec.kind == "title":
            draw.ellipse((120, 90, 205, 175), fill=f"#{GREEN}")
            draw.text((width // 2, 245), "Sati", fill=f"#{TEXT}", font=font_big, anchor="mm")
            draw.text((width // 2, 355), spec.kicker, fill=f"#{GREEN}", font=font_title, anchor="mm")
            draw.text((width // 2, 480), spec.body[0], fill=f"#{TEXT}", font=font_body, anchor="mm")
            draw.text((width // 2, 555), spec.body[1], fill=f"#{MUTED}", font=font_body, anchor="mm")
            draw.text((width // 2, 610), spec.body[2], fill=f"#{MUTED}", font=font_body, anchor="mm")
        elif spec.kind == "architecture":
            draw.text((90, 70), spec.title, fill=f"#{TEXT}", font=font_title)
            draw.rectangle((90, 140, 1510, 145), fill=f"#{GREEN}")
            draw.rounded_rectangle((105, 220, 1495, 650), radius=22, fill=f"#{CREAM}", outline="#C9DFBF", width=3)
            for row, line in enumerate(spec.body):
                draw.text((145, 300 + row * 68), line, fill=f"#{TEXT}", font=font_mono)
            draw.text((width // 2, 715), "UNO Q = edge hub: MCU bridge + Linux dashboard", fill=f"#{MUTED}", font=font_body, anchor="mm")
        else:
            draw.text((90, 70), spec.title, fill=f"#{TEXT}", font=font_title)
            draw.rectangle((90, 140, 1510, 145), fill=f"#{GREEN}")
        draw.text((1450, 835), f"{slide_number}/{len(SLIDES)}", fill=f"#{MUTED}", font=font_body)
        img.save(preview_dir / f"sati_pitch_slide_{slide_number}.png")


def verify_deck(pptx_path: Path) -> None:
    prs = Presentation(pptx_path)
    if len(prs.slides) != len(SLIDES):
        raise RuntimeError(f"Expected {len(SLIDES)} slides, found {len(prs.slides)}")

    body_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    if "[ทีมต้องเติม]" in body_text:
        raise RuntimeError("Placeholder [ทีมต้องเติม] found in slide body")

    with zipfile.ZipFile(pptx_path, "r") as package:
        note_names = sorted(name for name in package.namelist() if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name))
        if len(note_names) != len(SLIDES):
            raise RuntimeError(f"Expected {len(SLIDES)} notes slides, found {len(note_names)}")
        for name in note_names:
            if "<a:t>" not in package.read(name).decode("utf-8"):
                raise RuntimeError(f"Missing note text in {name}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the Sati hackathon pitch deck.")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--preview-dir", type=Path, default=None)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    build_deck(args.output)
    verify_deck(args.output)
    if args.preview_dir:
        render_preview(args.preview_dir, [1, 5])
    size_mb = args.output.stat().st_size / (1024 * 1024)
    print(f"wrote {args.output} ({size_mb:.2f} MB)")
    print(f"slides: {len(SLIDES)}")
    print(f"speaker notes: {len(SLIDES)}/{len(SLIDES)}")


if __name__ == "__main__":
    main()
